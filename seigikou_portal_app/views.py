from django.shortcuts import render,redirect
from django.urls import reverse_lazy
from django.urls import reverse
from django.http import HttpResponse,HttpResponseRedirect
from .forms import AccountForm, AddAccountForm, MemberForm # ユーザーアカウントフォーム
from django.contrib.auth import authenticate,login,logout
from django.contrib.auth.decorators import login_required
from django.conf import settings
from django.views.generic import (ListView,
                                  DetailView,
                                  CreateView,
                                  DeleteView,
                                  UpdateView,
                                  TemplateView)
from . import models
import os
import io
from pptx import Presentation
from pptx.util import Pt
import qrcode
from PIL import Image
from django.contrib import messages

'''
#ホーム画面
class Home(ListView):
    model = models.Event
    template_name = "Home.html"
'''
#一覧画面
class EventList(ListView):
    #Companyテーブル連携
    model = models.Event
    #レコード情報をテンプレートに渡すオブジェクト
    context_object_name = "event_list"
    #テンプレートファイル連携
    template_name = "Event_list.html"

#詳細画面
class EventDetail(DetailView):
    #Companyテーブル連携
    model = models.Event
    #レコード情報をテンプレートに渡すオブジェクト
    context_object_name = "event_detail"
    #テンプレートファイル連携
    template_name = "Event_detail.html"

#Create(イベント)画面
class EventCreateView(CreateView):
    #Companyテーブル連携
    model = models.Event
    #入力項目定義
    fields = ("day","starttime","endtime","location",
              "name","public","contents","question_url","member")
    #テンプレートファイル連携
    template_name = "Event_form.html"
    #更新後のリダイレクト先
    def get_success_url(self):
        return reverse('App:detail', kwargs={'pk': self.object.pk})

#メンバーリスト画面
class MemberList(ListView):
    #Companyテーブル連携
    model = models.Member
    #レコード情報をテンプレートに渡すオブジェクト
    context_object_name = "member_list"
    #テンプレートファイル連携
    template_name = "member_list.html"

#詳細画面
class MemberDetail(DetailView):
    #Companyテーブル連携
    model = models.Member
    #レコード情報をテンプレートに渡すオブジェクト
    context_object_name = "member_detail"
    #テンプレートファイル連携
    template_name = "Member_detail.html"

    def get(self, request, **kwargs):
        # アクティブユーザーでなければログインページ
        if not request.user.is_active:
            return redirect('/home')
        return super().get(request)

#Create(メンバー)画面
class MemberCreateView(CreateView):
    model = models.Member
    form_class = MemberForm
    template_name = "Member_form.html"
    success_url = reverse_lazy("App:memberlist")

#講演リスト画面
class KouenList(ListView):
    #Companyテーブル連携
    model = models.Kouen
    #レコード情報をテンプレートに渡すオブジェクト
    context_object_name = "kouen_list"
    #テンプレートファイル連携
    template_name = "kouen_list.html"

#詳細画面
class KouenDetail(DetailView):
    #Companyテーブル連携
    model = models.Kouen
    #レコード情報をテンプレートに渡すオブジェクト
    context_object_name = "kouen_detail"
    #テンプレートファイル連携
    template_name = "kouen_detail.html"

#Create(講演)画面
class KouenCreateView(CreateView):
    #Companyテーブル連携
    model = models.Kouen
    #入力項目定義
    fields = ("koushi","name","youshi","event")
    #テンプレートファイル連携
    template_name = "Event_form.html"
    #作成後のリダイレクト先
    success_url = reverse_lazy("App:kouenlist")

#Upadate画面(イベント情報)
class EventUpdateView(UpdateView):
    #入力項目定義
    fields =  ("day","starttime","endtime","location",
              "name","public","contents","question_url","member")
    #Companyテーブル連携
    model = models.Event
    #テンプレートファイル連携
    template_name = "Event_form.html"
    #更新後のリダイレクト先
    def get_success_url(self):
        return reverse('App:detail', kwargs={'pk': self.object.pk})

#更新画面(メンバー情報)
class MemberUpdateView(UpdateView):
    model = models.Member
    form_class = MemberForm
    template_name = "Member_form.html"
    success_url = reverse_lazy("App:memberlist")
    
#更新(講演)画面
class KouenUpdateView(UpdateView):
    #入力項目定義
    fields = ("koushi","name","youshi","event")
    #テーブル連携
    model = models.Kouen
    #テンプレートファイル連携
    template_name = "Event_form.html"
    #作成後のリダイレクト先
    success_url = reverse_lazy("App:kouenlist")

#削除画面
class EventDeleteView(DeleteView):
    #イベントテーブル連携
    model = models.Event
    #テンプレートファイル連携
    template_name = "Event_delete.html"
    #削除後のリダイレクト先
    success_url = reverse_lazy("App:eventlist")

    def get(self, request, **kwargs):
        if not request.user.is_staff:
            return redirect('/eventlist')
        return super().get(request)


#ログイン
def Login(request):
    # POST
    if request.method == 'POST':
        # フォーム入力のユーザーID・パスワード取得
        ID = request.POST.get('userid')
        Pass = request.POST.get('password')

        # Djangoの認証機能
        user = authenticate(username=ID, password=Pass)

        # ユーザー認証
        if user:
            #ユーザーアクティベート判定
            if user.is_active:
                # ログイン
                login(request,user)
                # ホームページ遷移
                return HttpResponseRedirect(reverse('App:home'))
            else:
                # アカウント利用不可
                return HttpResponse("アカウントが有効ではありません")
        # ユーザー認証失敗
        else:
            return HttpResponse("ログインIDまたはパスワードが間違っています")
    # GET
    else:
        return render(request, 'login.html')


#ログアウト
@login_required
def Logout(request):
    logout(request)
    # ログイン画面遷移
    return HttpResponseRedirect(reverse('App:Login'))


#ホーム
@login_required
def home(request):
    params = {"UserID":request.user,}
    return render(request, "home.html",context=params)


#新規登録
class  AccountRegistration(TemplateView):

    def __init__(self):
        self.params = {
        "AccountCreate":False,
        "account_form": AccountForm(),
        "add_account_form":AddAccountForm(),
        }

    #Get処理
    def get(self,request):
        self.params["account_form"] = AccountForm()
        self.params["add_account_form"] = AddAccountForm()
        self.params["AccountCreate"] = False
        return render(request,"register.html",context=self.params)

    #Post処理
    def post(self,request):
        self.params["account_form"] = AccountForm(data=request.POST)
        self.params["add_account_form"] = AddAccountForm(data=request.POST)

        #フォーム入力の有効検証
        if self.params["account_form"].is_valid() and self.params["add_account_form"].is_valid():
        
            # キーワードの検証
            secret_word = self.params["add_account_form"].cleaned_data['secret_word']
            # ここで正しいキーワードを設定してください
            correct_secret_word = os.environ.get('SECRET_WORD') #ココが秘密のキーワード

            if secret_word == correct_secret_word:
                # アカウント情報をDB保存
                account = self.params["account_form"].save()
                # パスワードをハッシュ化
                account.set_password(account.password)
                # ハッシュ化パスワード更新
                account.save()

                # 下記追加情報
                # 下記操作のため、コミットなし
                add_account = self.params["add_account_form"].save(commit=False)
                # AccountForm & AddAccountForm 1vs1 紐付け
                add_account.user = account

                # 画像アップロード有無検証
                if 'account_image' in request.FILES:
                    add_account.account_image = request.FILES['account_image']

                # モデル保存
                add_account.save()

                # アカウント作成情報更新
                self.params["AccountCreate"] = True

            else:
                # キーワードが正しくない場合はエラーメッセージを表示
                 messages.error(request, '秘密の言葉が違うよ')

        else:
            # フォームが有効でない場合
            print(self.params["account_form"].errors)
            messages.error(request, '入力内容にエラーがあります')

        return render(request,"register.html",context=self.params)
    

# PPT出力準備（次回イベント選択など）
class EventPpt(DetailView):
    #Eventテーブル連携
    model = models.Event
    #レコード情報をテンプレートに渡すオブジェクト
    context_object_name = "event_detail"
    #テンプレートファイル連携
    template_name = "Event_ppt.html"

    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        context['events'] = models.Event.objects.all()  # すべての著者を追加
        return context


def EventPptDownload2(request, pk):
    if request.method == 'POST':
        next_event_pk = request.POST.get('item')

        # ダウンロードするファイルのパス
        file_path = os.path.join(settings.MEDIA_ROOT, 'template.pptx')
        # **kwargs引数に指定されたpkのレコードを取り出す
        event = models.Event.objects.get(pk=pk)
        nextevent = models.Event.objects.get(pk=next_event_pk)
        print(event.id, event.name)

        prs = Presentation(file_path)

        for slide in prs.slides:
            for shape in slide.shapes:

                # P.1のテーブル書き換え
                if shape.has_text_frame and shape.name == "txtbox_reikai_contents":
                    # 例会内容
                    shape.text_frame.text = event.contents
                    # 改行でパラグラフが複数あると書式が変わらない
                    for para in shape.text_frame.paragraphs:
                        para.font.size = Pt(24)
                        for run in para.runs:
                            # テキスト内の "_x000D_" を削除する ※改行(¥n)は存在するので削除のみで可
                            run.text = run.text.replace("_x000D_", "")
                elif shape.has_text_frame and shape.name == "txtbox_koushi_endai":
                    # 講師・演題
                    if event.kouen:
                        shape.text_frame.text = f'''演題：{event.kouen.name}
講師：{event.kouen.koushi}({event.kouen.koushi.company})'''
                    else:
                        # kouenフィールドに値がない場合
                        shape.text_frame.text = f'''未定'''
                    # 
                    # 改行でパラグラフが複数あると書式が変わらない
                    for para in shape.text_frame.paragraphs:
                        para.font.size = Pt(24)

                # P.2の書き換え
                elif shape.name == "txtbox_next_title":
                    shape.text_frame.text = nextevent.name
                    
                elif shape.name == "txtbox_next_schedule":
                    shape.text_frame.text = f'''{str(nextevent.day)}
{nextevent.starttime.strftime("%H:%M")}〜
{nextevent.endtime.strftime("%H:%M")}'''
                    
                elif shape.name == "txtbox_next_koushi_endai":
                    if nextevent.kouen:
                        # 講師・演題
                        shape.text_frame.text = f'''演題：{nextevent.kouen.name}
講師：{nextevent.kouen.koushi}({nextevent.kouen.koushi.company})'''
                    else:
                        shape.text_frame.text = f'''未定'''

                    # 改行でパラグラフが複数あると書式が変わらない
                    for para in shape.text_frame.paragraphs:
                        para.font.size = Pt(24)
                    
                elif shape.name == "txtbox_ enquete_url":
                    shape.text_frame.text = f'''{event.question_url}'''
                    shape.text_frame.paragraphs[0].font.size = Pt(10)
        
        # QRコードの書き換え
        img = qrcode.make(event.question_url)
        pic_io = io.BytesIO()
        img.save(pic_io, format="PNG")
        pic_io.seek(0)
        
        # pngとして保存したQRコードをPPTに貼り付け
        # リファレンス https://python-pptx.readthedocs.io/en/latest/api/shapes.html#picture-objects
        prs.slides[1].shapes.add_picture(pic_io, 9469796, 1147979, 1489415)


        # 絵の位置を調べる
        # imgshp = prs.slides[1].shapes[6]
        # print(f'top:{imgshp.top}, left:{imgshp.left}, width:{imgshp.width}, height:{imgshp.height}')
        # top:1147979, left:9469796, width:1489415, height:1489415


        
        
        # ファイルをバイナリストリームとして保存
        ppt_io = io.BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)  # ストリームの先頭に戻る

        # レスポンスの作成
        response = HttpResponse(ppt_io.getvalue(), content_type='application/vnd.ms-powerpoint')
        response['Content-Disposition'] = 'attachment; filename="reikai.pptx"'
        response['Cache-Control'] = 'no-cache, no-store, must-revalidate'  # キャッシュを無効にする
        response['Pragma'] = 'no-cache'  # 古いブラウザ用
        response['Expires'] = '0'  # プロキシサーバーのキャッシュを防ぐ
        return response #redirect('App:detail', pk=pk)

    else:
        return render(request, 'Event_list.html')


#PPT出力
#関数バージョンに移行したので不要。タイミング見て消す
# class EventPptDownload(DetailView):
#     #Companyテーブル連携
#     model = models.Event
#     #レコード情報をテンプレートに渡すオブジェクト
#     context_object_name = "event_detail"
#     #テンプレートファイル連携
#     template_name = "event_detail.html"

#     def get(self, request, *args, **kwargs):
#         # ダウンロードするファイルのパス
#         file_path = os.path.join(settings.MEDIA_ROOT, 'template.pptx')
        
#         # kwargsの中身確認。ディクショナリ形式でpkにイベントIDが格納
#         # print("<print **kwargs>")
#         # for key, value in kwargs.items():
#         #     print(f"{key}: {value}")

#         # 全件取得
#         # events = models.Event.objects.all()
#         # for event in events:
#         #     print(event.id, event.name)

#         # **kwargs引数に指定されたpkのレコードを取り出す
#         event = models.Event.objects.get(pk = kwargs["pk"])
#         print(event.id, event.name)


#         prs = Presentation(file_path)

#         # P.1のテーブル書き換え
#         for slide in prs.slides:
#             for shape in slide.shapes:
#                 if shape.has_text_frame and shape.name == "txtbox_reikai_contents":
#                     # 例会内容
#                     shape.text_frame.text = event.contents
#                     # 改行でパラグラフが複数あると書式が変わらない
#                     for para in shape.text_frame.paragraphs:
#                         para.font.size = Pt(24)
#                         for run in para.runs:
#                             # テキスト内の "_x000D_" を削除する ※改行(¥n)は存在するので削除のみで可
#                             run.text = run.text.replace("_x000D_", "")
#                 elif shape.has_text_frame and shape.name == "txtbox_koushi_endai":
#                     # 講師・演題
#                     shape.text_frame.text = f'''演題：{event.kouen.name}
# 講師：{event.kouen.koushi}({event.kouen.koushi.company})'''
#                     # 
#                     # 改行でパラグラフが複数あると書式が変わらない
#                     for para in shape.text_frame.paragraphs:
#                         para.font.size = Pt(24)

#         """
#         table.cell(2,2).text = f'''講師：{kouen_dict[kouen_id][3].value}
#         演題：{kouen_dict[kouen_id][4].value}'''
#         # 改行でパラグラフが複数あると書式が変わらない
#         for para in table.cell(2,2).text_frame.paragraphs:
#             para.font.size = Pt(24)

#         table.cell(5,2).text = f'''青年技術士交流会 例会
#         {event_list[en][7].value}'''
#         table.cell(5,2).text_frame.paragraphs[0].font.size = Pt(24)


#         # 2ページ目Shapesの中身を確認※調査用
#         for shp in prs.slides[1].shapes:
#             print(f'{type(shp)}, name:{shp.name}')

#         # P.2のテーブル書き換え
#         table = prs.slides[1].shapes[3].table

#         # 次回イベント
#         table.cell(0,0).text = f'{event_list[nxten][5].value}'
#         table.cell(0,0).text_frame.paragraphs[0].font.size = Pt(24)
#         # 次回開催日時
#         day_fmt = "%m/%d %a"
#         tm_fmt = "%H:%M"
#         table.cell(0,1).text = f'''{event_list[nxten][1].value.strftime(day_fmt)}
#         {event_list[nxten][2].value.strftime(tm_fmt)}～'''
#         table.cell(0,1).text_frame.paragraphs[0].font.size = Pt(24)
#         table.cell(0,1).text_frame.paragraphs[1].font.size = Pt(24)
#         # 次回講演内容
#         table.cell(0,2).text = f'''講師：{kouen_dict[nxtkouen_id][3].value}
#         演題：{kouen_dict[nxtkouen_id][4].value}
#         '''
#         if event_list[nxten][6].value == '青年ML':
#             table.cell(0,2).text += "※青年技術士交流会の内部勉強会になります。"

#         table.cell(0,2).text_frame.paragraphs[0].font.size = Pt(24)
#         table.cell(0,2).text_frame.paragraphs[1].font.size = Pt(24)

#         print(type(event_list[nxten][1].value))

#         # table = prs.slides[1].shapes[1].table

#         # table.cell(2,2).text = f'''講師：{kouen_dict[kouen_id][3].value}
#         # 演題：{kouen_dict[kouen_id][4].value}'''


#         # QRコードの作成とPPTへの貼り付け
#         # 参考）https://tat-pytone.hatenablog.com/entry/2019/06/22/153657
#         # ------------------------------------------------------------------

#         # QRコード生成
#         url = event_list[en][11].value
#         img = qrcode.make(url)
#         img.save('qrcode.png')

#         # テキストボックスのURL書き換え
#         urltxtbox = prs.slides[1].shapes[7].text_frame.paragraphs[0]
#         urltxtbox.text = url
#         urltxtbox.font.size = Pt(12)
#         urltxtbox.font.underline = True

#         # pngとして保存したQRコードをPPTに貼り付け
#         # リファレンス https://python-pptx.readthedocs.io/en/latest/api/shapes.html#picture-objects
#         prs.slides[1].shapes.add_picture("qrcode.png", 9469796, 1147979, 1489415)


#         # 絵の位置を調べる
#         # imgshp = prs.slides[1].shapes[6]
#         # print(f'top:{imgshp.top}, left:{imgshp.left}, width:{imgshp.width}, height:{imgshp.height}')
#         # top:1147979, left:9469796, width:1489415, height:1489415
#         """


#         # ファイルをバイナリストリームとして保存
#         ppt_io = io.BytesIO()
#         prs.save(ppt_io)
#         ppt_io.seek(0)  # ストリームの先頭に戻る

#         # レスポンスの作成
#         response = HttpResponse(ppt_io.getvalue(), content_type='application/vnd.ms-powerpoint')
#         response['Content-Disposition'] = 'attachment; filename="reikai.pptx"'
#         response['Cache-Control'] = 'no-cache, no-store, must-revalidate'  # キャッシュを無効にする
#         response['Pragma'] = 'no-cache'  # 古いブラウザ用
#         response['Expires'] = '0'  # プロキシサーバーのキャッシュを防ぐ

#         # # 別名でPPT保存して終了
#         # # ------------------------------------------------------------------
#         # prs.save('test.pptx')

#         # # ファイルをバイナリ形式で読み込む
#         # with open(file_path, 'rb') as f:
#         #     file_data = f.read()
        
#         # # HttpResponseオブジェクトを作成して、ファイルをダウンロードさせる
#         # response = HttpResponse(file_data, content_type='application/vnd.ms-powerpoint')
#         # response['Content-Disposition'] = 'attachment; filename="test.pptx"'
        

#         return response

